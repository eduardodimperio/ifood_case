"""
Script para gerar a apresentação de 5 slides do iFood Case.
Execução: python3 presentation/create_slides.py
Saída:    presentation/ifood_case_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os, textwrap

# ─── Paleta de Cores ────────────────────────────────────────────────────────
IFOOD_RED    = RGBColor(0xEA, 0x1D, 0x2C)   # vermelho iFood
DARK_GRAY    = RGBColor(0x2C, 0x2C, 0x2C)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE  = RGBColor(0x34, 0x98, 0xDB)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
MID_GRAY     = RGBColor(0x95, 0xA5, 0xA6)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── Helpers ────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=DARK_GRAY, bullet="●  "):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    kwargs = {"left": Inches(left), "top": Inches(top)}
    if width:  kwargs["width"]  = Inches(width)
    if height: kwargs["height"] = Inches(height)
    slide.shapes.add_picture(path, **kwargs)

def add_header_bar(slide, title, subtitle=None):
    """Barra vermelha de título no topo."""
    add_rect(slide, 0, 0, 13.33, 1.3, IFOOD_RED)
    add_text(slide, title, 0.4, 0.1, 12, 0.75,
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.78, 12, 0.5,
                 font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 7.2, 13.33, 0.3, DARK_GRAY)
    add_text(slide, "iFood — Case Técnico Data Science | Confidencial", 0.3, 7.21, 13, 0.28,
             font_size=9, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ─── Criação do PPTX ────────────────────────────────────────────────────────
PRESENTATION_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT     = os.path.dirname(PRESENTATION_DIR)
IMG = PRESENTATION_DIR

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completamente em branco

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA / O DESAFIO
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(sl, 0, 0, 5.2, 7.5, IFOOD_RED)

add_text(sl, "iFood", 0.45, 0.6, 4.5, 1.0,
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Distribuição Inteligente\nde Cupons e Ofertas", 0.45, 1.55, 4.5, 2.0,
         font_size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Case Técnico · Data Science", 0.45, 3.4, 4.5, 0.6,
         font_size=13, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=True)
add_text(sl, "Eduardo Dimperio", 0.45, 6.7, 4.5, 0.5,
         font_size=11, color=WHITE, align=PP_ALIGN.LEFT)

# Lado direito: o desafio
add_text(sl, "O Desafio", 5.8, 0.8, 7, 0.7,
         font_size=26, bold=True, color=IFOOD_RED, align=PP_ALIGN.LEFT)

challenges = [
    "17.000 clientes com perfis muito distintos",
    "10 ofertas de tipos diferentes (BOGO, desconto, informacional)",
    "Múltiplos canais de marketing (e-mail, mobile, social, web)",
    "Hoje: distribuição aleatória ou regras manuais",
    "Oportunidade: distribuição personalizada = mais conversões com menos custo",
]
add_bullet_list(sl, challenges, 5.8, 1.65, 7.1, 4.5, font_size=15, color=DARK_GRAY)

add_text(sl, "📊 Dados disponíveis: 300k eventos · 17k clientes · histórico de 30 dias",
         5.8, 6.5, 7.1, 0.7,
         font_size=12, italic=True, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — O QUE OS DADOS REVELAM
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_header_bar(sl, "O Que os Dados Revelam",
               "Análise exploratória: padrões de conversão identificados")

# Imagem EDA
add_image(sl, os.path.join(IMG, "eda_conversion_rates.png"), 0.3, 1.4, 8.0)

# Insights à direita
add_text(sl, "Insights-Chave", 8.6, 1.35, 4.4, 0.5,
         font_size=17, bold=True, color=DARK_GRAY)

insights = [
    "Ofertas informacionais convertem +28% mais que BOGO",
    "Canal social tem melhor taxa p/ descontos",
    "Clientes 40–60 anos são mais responsivos",
    "Ticket médio: R$ 12 — limite mínimo das ofertas compatível",
    "Perfis incompletos (~30%) têm comportamento distinto",
]
add_bullet_list(sl, insights, 8.6, 1.9, 4.4, 4.5, font_size=13, color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — A SOLUÇÃO: MODELO DE RECOMENDAÇÃO
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_header_bar(sl, "A Solução: Modelo de Recomendação Personalizado",
               "Para cada cliente → oferta com maior lucro esperado")

# Caixas de processo (pipeline)
boxes = [
    ("1. Dados", "Perfil do cliente\nHistórico de compras\nMetadados das ofertas", ACCENT_BLUE),
    ("2. Modelo", "LightGBM\nClassificação binária\nAUC = 0.92", IFOOD_RED),
    ("3. Score", "P(conversão)\npor par\n(cliente × oferta)", ACCENT_GREEN),
    ("4. Decisão", "Net EV = P × (ticket − desconto)\nOferta com maior lucro esperado", DARK_GRAY),
]
box_w, box_h = 2.8, 2.5
start_x = 0.45
for i, (title, body, color) in enumerate(boxes):
    x = start_x + i * (box_w + 0.25)
    add_rect(sl, x, 1.5, box_w, box_h, color)
    add_text(sl, title, x + 0.15, 1.6, box_w - 0.3, 0.55,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + 0.1, 2.2, box_w - 0.2, box_h - 0.8, WHITE)
    add_text(sl, body, x + 0.2, 2.3, box_w - 0.4, box_h - 1.0,
             font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(sl, "→", x + box_w + 0.0, 2.45, 0.28, 0.5,
                 font_size=22, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Features mais importantes
add_text(sl, "Principais sinais preditivos identificados:",
         0.45, 4.25, 12.5, 0.5, font_size=14, bold=True, color=DARK_GRAY)
add_image(sl, os.path.join(IMG, "feature_importance.png"), 0.45, 4.6, 6.5)

note = (
    "O modelo aprende: clientes com alto ticket médio preferem BOGO; "
    "clientes mais sensíveis ao preço respondem melhor a descontos menores. "
    "Isso permite enviar a oferta certa para cada perfil."
)
add_text(sl, note, 7.2, 4.7, 5.9, 2.3, font_size=12,
         italic=True, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS E IMPACTO NO NEGÓCIO
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_header_bar(sl, "Resultados: Impacto no Negócio",
               "Comparativo entre estratégias de distribuição de ofertas")

add_image(sl, os.path.join(IMG, "business_impact.png"), 0.3, 1.4, 8.5)

# KPIs à direita
add_text(sl, "Performance do Modelo", 9.1, 1.4, 4.0, 0.5,
         font_size=16, bold=True, color=DARK_GRAY)

kpis = [
    ("AUC = 0.92", "Alta capacidade de discriminação", ACCENT_GREEN),
    ("+16%", "de conversão vs distribuição aleatória", ACCENT_BLUE),
    ("+67%", "de resultado líquido vs aleatório", IFOOD_RED),
    ("R$ 45k", "de incremento estimado* por ciclo", DARK_GRAY),
]
y = 2.0
for val, label, color in kpis:
    add_rect(sl, 9.1, y, 4.0, 0.9, color)
    add_text(sl, val, 9.15, y + 0.03, 1.5, 0.5,
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, label, 10.65, y + 0.15, 2.4, 0.65,
             font_size=11, color=WHITE, align=PP_ALIGN.LEFT)
    y += 1.05

add_text(sl, "* Estimativa baseada em ticket médio histórico de R$12. Validação via A/B test recomendada.",
         0.3, 6.85, 12.7, 0.35,
         font_size=9, italic=True, color=MID_GRAY, align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PRÓXIMOS PASSOS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_header_bar(sl, "Próximos Passos & Roadmap",
               "Da prova de conceito à implementação em produção")

steps = [
    (
        "Curto prazo (1–4 semanas)",
        IFOOD_RED,
        [
            "Teste A/B controlado: grupo controle (aleatório) vs grupo modelo",
            "Definição da métrica de sucesso: GMV incremental por oferta enviada",
            "Validação do lift causal (não apenas correlação)",
        ]
    ),
    (
        "Médio prazo (1–3 meses)",
        ACCENT_BLUE,
        [
            "Correção do temporal leakage: features calculadas antes do recebimento da oferta",
            "Retraining semanal com dados frescos (dados novos chegam continuamente)",
            "Explorar bandits contextuais (UCB, Thompson Sampling) para exploração/exploitação",
        ]
    ),
    (
        "Longo prazo (3–6 meses)",
        ACCENT_GREEN,
        [
            "Pipeline real-time: scoring on-demand ao enviar oferta",
            "Otimização de budget: alocação de orçamento de desconto por segmento",
            "Modelo causal (uplift modeling) para identificar clientes genuinamente incremetais",
        ]
    ),
]

x = 0.35
for title, color, bullets in steps:
    add_rect(sl, x, 1.45, 4.1, 5.35, color)
    add_text(sl, title, x + 0.1, 1.52, 3.9, 0.65,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + 0.08, 2.18, 3.94, 4.57, WHITE)
    add_bullet_list(sl, bullets, x + 0.18, 2.28, 3.75, 4.4,
                    font_size=12, color=DARK_GRAY, bullet="▸  ")
    x += 4.35

add_text(sl, "\"A verdadeira vantagem competitiva está em aprender continuamente "
             "o que funciona para cada cliente — e agir mais rápido que o mercado.\"",
         0.6, 6.8, 12.1, 0.55,
         font_size=11, italic=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# ─── Salvar ─────────────────────────────────────────────────────────────────
out = os.path.join(PRESENTATION_DIR, "ifood_case_presentation.pptx")
prs.save(out)
print(f"Apresentação salva em: {out}")
